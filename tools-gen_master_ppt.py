# -*- coding: utf-8 -*-
# NICE 售前 master V3 — 亮底版式 + 渐变/投影/图标/主视觉/数据图/实拍
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
from PIL import Image

NAVY =RGBColor(0x0A,0x1A,0x3D); HBAND=RGBColor(0x12,0x2E,0x63); DARK=RGBColor(0x0B,0x10,0x24)
BLUE =RGBColor(0x24,0x6B,0xFF); CYAN=RGBColor(0x00,0xB5,0xE3); MAG=RGBColor(0x7D,0x3C,0xFF)
VIO  =RGBColor(0x5A,0x1B,0xDF); LILAC=RGBColor(0xC8,0x7D,0xFF)
WHITE=RGBColor(0xFF,0xFF,0xFF); FOG=RGBColor(0xD0,0xD5,0xE2)
TITLED=RGBColor(0x0A,0x1A,0x3D); BODY=RGBColor(0x5B,0x67,0x82); MUTE=RGBColor(0x95,0x9F,0xB5)
CLINE=RGBColor(0xE2,0xE8,0xF3); PANEL=RGBColor(0xF5,0xF8,0xFD); ENSUB=RGBColor(0x8F,0xB2,0xF0)
CN='微软雅黑'; EN='Arial'; ENB='Arial Black'
CT=PP_ALIGN.CENTER; LF=PP_ALIGN.LEFT; RT=PP_ALIGN.RIGHT; MID=MSO_ANCHOR.MIDDLE; TOP=MSO_ANCHOR.TOP
def hx(c): return '%02X%02X%02X'%(c[0],c[1],c[2])

AS='/tmp/assets/white/w-%s.png'
ICN={'search':'8-2','palette':'8-3','cube':'8-4','touch':'8-5','factory':'8-6','bulb':'15-9',
 'gear':'14-13','gears':'24-4','chip':'13-4','code':'14-6','org':'24-2','eye':'16-3','finger':'16-2',
 'brain':'17-4','flash':'17-2','leaf':'17-7','seat':'24-5','car':'19-3','magnet':'19-4','cup':'19-5',
 'trophy':'22-2','robot':'25-3','hand':'26-2','check':'30-2','target':'14-15','clip':'14-2','bench':'15-7','wheel':'14-15'}
LOC=['7-2','7-3','7-4','7-5','7-6']

# 预裁船舶驾驶舱图 → 2.7:1
def prep_ship():
    src='/tmp/assets/ship_image43.jpeg'
    if not os.path.exists(src): src='/tmp/assets/ship_image40.jpeg'
    im=Image.open(src).convert('RGB'); w,h=im.size; ar=2.7
    nw=min(w,int(h*ar)); nh=int(nw/ar); l=(w-nw)//2; t=(h-nh)//2
    im.crop((l,t,l+nw,t+nh)).save('/tmp/assets/ship_crop.jpg',quality=88)
prep_ship()

def crop_ar(src,dst,ar,q=90):
    if not os.path.exists(src): return False
    im=Image.open(src).convert('RGB'); w,h=im.size
    if w/h>ar: nw=int(h*ar); nh=h
    else: nw=w; nh=int(w/ar)
    l=(w-nw)//2; t=(h-nh)//2; im.crop((l,t,l+nw,t+nh)).save(dst,quality=q); return True
# iX-2024 红点座舱 → NICE概念页(横构图)
crop_ar('/tmp/pdfimg/jomec_p37_576.jpeg','/tmp/assets/ix.jpg',3.7/2.05)
# JoySpace+ 座舱 → 案例页横幅
crop_ar('/tmp/pdfimg/jomec_p36_572.jpeg','/tmp/assets/joy.jpg',4.1/1.45)
import glob as _glob
def gpath(k):
    r=_glob.glob('/tmp/pdfimg/%s.*'%k); return r[0] if r else ''
crop_ar('/tmp/pdfimg/jomec_p37_576.jpeg','/tmp/assets/cover.jpg',3.9/2.7)   # iX红点座舱→封面hero
crop_ar(gpath('jomec_p33_553'),'/tmp/assets/bench.jpg',2.6/1.0)        # 测试台架实拍
crop_ar(gpath('jomec_p29_504'),'/tmp/assets/seat.jpg',2.6/1.0)         # 座椅/机构渲染
crop_ar(gpath('jomec_p26_430'),'/tmp/assets/ctrl.jpg',2.6/1.0)        # 代码/软硬件
def pic(s,x,y,w,h,path,border=None,shd=True):
    if not os.path.exists(path): return None
    p=s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if border: p.line.color.rgb=border; p.line.width=Pt(1.2)
    if shd: shadow(p,alpha=78000)
    return p

prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
BL=prs.slide_layouts[6]

def _f(r,nm,sz,bd,cl,it=False):
    r.font.size=Pt(sz); r.font.bold=bd; r.font.italic=it; r.font.color.rgb=cl; r.font.name=nm
    rPr=r._r.get_or_add_rPr()
    for tg in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(tg))
        if e is None: e=rPr.makeelement(qn(tg),{}); rPr.append(e)
        e.set('typeface',nm)

def grad(sp,c1,c2,ang=2700000):
    spPr=sp._element.spPr
    for tg in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e=spPr.find(qn(tg));
        if e is not None: spPr.remove(e)
    g=spPr.makeelement(qn('a:gradFill'),{}); lst=g.makeelement(qn('a:gsLst'),{})
    for pos,col in ((0,c1),(100000,c2)):
        gs=g.makeelement(qn('a:gs'),{'pos':str(pos)}); sc=g.makeelement(qn('a:srgbClr'),{'val':hx(col)})
        gs.append(sc); lst.append(gs)
    g.append(lst); g.append(g.makeelement(qn('a:lin'),{'ang':str(ang),'scaled':'1'}))
    ln=spPr.find(qn('a:ln'))
    (ln.addprevious(g) if ln is not None else spPr.append(g))

def shadow(sp,blur=80000,dist=30000,dr=5400000,alpha=70000,col='17244A'):
    spPr=sp._element.spPr; el=spPr.find(qn('a:effectLst'))
    if el is None: el=spPr.makeelement(qn('a:effectLst'),{}); spPr.append(el)
    sh=el.makeelement(qn('a:outerShdw'),{'blurRad':str(blur),'dist':str(dist),'dir':str(dr),'rotWithShape':'0'})
    c=el.makeelement(qn('a:srgbClr'),{'val':col}); c.append(el.makeelement(qn('a:alpha'),{'val':str(alpha)}))
    sh.append(c); el.append(sh)

def slide(bg):
    s=prs.slides.add_slide(BL)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s,r

def rect(s,x,y,w,h,fill=None,line=None,lw=0.75,sh=MSO_SHAPE.RECTANGLE,adj=None):
    sp=s.shapes.add_shape(sh,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if adj is not None:
        try: sp.adjustments[0]=adj
        except: pass
    return sp

def txt(s,x,y,w,h,paras,al=LF,an=TOP,after=2,ls=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=an
    tf.margin_left=Inches(0.03); tf.margin_right=Inches(0.03); tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01)
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        pa.alignment=al; pa.space_after=Pt(after); pa.space_before=Pt(0); pa.line_spacing=ls
        for seg in p:
            t,nm,sz,bd,cl=seg[:5]; it=seg[5] if len(seg)>5 else False
            r=pa.add_run(); r.text=t; _f(r,nm,sz,bd,cl,it)
    return tb

def circ(s,cx,cy,d,fill,fill2=None,label='',icon=None,lcolor=WHITE,lsize=13):
    o=rect(s,cx-d/2,cy-d/2,d,d,fill=fill,sh=MSO_SHAPE.OVAL)
    if fill2: grad(o,fill,fill2,2700000)
    if icon and os.path.exists(AS%ICN.get(icon,'')):
        sz=d*0.52; s.shapes.add_picture(AS%ICN[icon],Inches(cx-sz/2),Inches(cy-sz/2),Inches(sz),Inches(sz))
    elif label: txt(s,cx-d/2,cy-d/2,d,d,[[(label,EN,lsize,True,lcolor)]],al=CT,an=MID)
    return o

def chead(s,no,zh,en,accent=BLUE):
    b=rect(s,0,0,10,0.92,fill=HBAND); grad(b,NAVY,RGBColor(0x1B,0x44,0x8E),0)
    rect(s,0,0,0.92,0.92,fill=NAVY)
    rect(s,0.92,0,0.045,0.92,fill=accent)
    txt(s,0,0,0.92,0.92,[[(f'{no:02d}',EN,24,True,WHITE)]],al=CT,an=MID)
    txt(s,1.18,0.16,7.4,0.45,[[(zh,CN,17,True,WHITE)]])
    txt(s,1.2,0.6,7.4,0.25,[[(en,EN,9,True,ENSUB)]])
    for r_ in range(3):
        for c_ in range(6):
            rect(s,8.55+c_*0.16,0.22+r_*0.17,0.045,0.045,fill=RGBColor(0x3A,0x5C,0x9E),sh=MSO_SHAPE.OVAL)
    rect(s,0,0.915,10,0.03,fill=accent)

def ccard(s,x,y,w,h,topbar=BLUE):
    c=rect(s,x,y,w,h,fill=WHITE,line=CLINE,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05); shadow(c)
    rect(s,x+0.12,y,w-0.24,0.05,fill=topbar)
    return c

def story(s,text,label='微故事',accent=MAG,x=0.55,y=1.08,w=8.9,h=0.6):
    c=rect(s,x,y,w,h,fill=PANEL,line=CLINE,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12); shadow(c,blur=50000,dist=18000,alpha=80000)
    rect(s,x,y+0.08,0.06,h-0.16,fill=accent)
    txt(s,x+0.22,y,w-0.4,h,[[(label+' · ',CN,10,True,accent),(text,CN,10,False,BODY)]],an=MID,ls=1.12)

def bottom(s,bold,rest='',x=0.55,y=4.36,w=8.9,h=0.64):
    c=rect(s,x,y,w,h,fill=NAVY,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.12); grad(c,NAVY,RGBColor(0x29,0x1C,0x5E),0); shadow(c,alpha=78000)
    rect(s,x+0.0,y+0.14,0.06,h-0.28,fill=CYAN)
    txt(s,x+0.32,y,w-0.6,h,[[(bold,CN,12.5,True,WHITE),('  '+rest if rest else '',CN,11,False,FOG)]],an=MID,ls=1.1)

def cfoot(s,n):
    rect(s,0.55,5.3,8.9,0.012,fill=CLINE)
    txt(s,0.9,5.33,7,0.24,[[('卓迈 JOMEC',CN,8,True,RGBColor(0x6A,0x78,0x96)),('  ·  NICE 智能座舱体验舱 · 把硬件优势变成体验话语权',CN,8,False,MUTE)]])
    txt(s,8.7,5.33,0.75,0.24,[[(f'{n:02d}',EN,8,True,RGBColor(0x6A,0x78,0x96))]],al=RT)

def techlines(s,bg):
    # 几何科技线条(暗底主视觉)
    for i,(x,y,w,h,c,a) in enumerate([(6.3,-0.5,4.5,4.5,BLUE,1),(7.6,1.4,3.6,3.6,VIO,1)]):
        o=rect(s,x,y,w,h,sh=MSO_SHAPE.OVAL); o.fill.background(); o.line.color.rgb=c; o.line.width=Pt(1)
    rect(s,0,5.18,10,0.012,fill=RGBColor(0x24,0x3A,0x70))
    for gx in range(6):
        rect(s,6.2+gx*0.62,0.0,0.012,5.625,fill=RGBColor(0x16,0x28,0x55))

def divider(part,zh,en,quote,big):
    s,bgr=slide(DARK); grad(bgr,DARK,RGBColor(0x16,0x12,0x3A),2700000)
    techlines(s,DARK)
    txt(s,6.0,0.3,4.0,3.6,[[(big,ENB,200,True,RGBColor(0x16,0x24,0x4E))]],al=RT)
    txt(s,0.6,0.85,3,0.4,[[(part,EN,12,True,MAG)]])
    rect(s,0.62,1.6,0.08,1.5,fill=MAG)
    txt(s,0.95,1.5,8.0,0.5,[[(en,EN,12,True,CYAN)]])
    txt(s,0.93,1.95,8.4,0.9,[[(zh,CN,30,True,WHITE)]])
    q=rect(s,0.95,3.25,7.6,1.0,fill=RGBColor(0x10,0x22,0x4A),line=RGBColor(0x2B,0x3A,0x66),sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08); shadow(q,alpha=75000)
    rect(s,0.95,3.33,0.06,0.84,fill=CYAN)
    txt(s,1.22,3.25,7.2,1.0,[[(quote,CN,12,False,FOG)]],an=MID,ls=1.2)
    return s

C4=[0.55,2.83,5.11,7.39]; W4=2.12; C3=[0.55,3.58,6.61]; W3=2.9; C2=[0.55,5.05]; W2=4.4
P=[0]
def pg(): P[0]+=1; return P[0]

# ============ 1 封面(主视觉)
s,bgr=slide(NAVY); grad(bgr,RGBColor(0x07,0x14,0x33),RGBColor(0x1A,0x12,0x40),2700000)
techlines(s,NAVY)
rect(s,0,0,10,0.09,fill=MAG); rect(s,0,5.535,10,0.09,fill=BLUE)
rect(s,0.7,1.08,0.08,0.5,fill=CYAN)
txt(s,0.96,1.05,8,0.4,[[('JOMEC ｜ 卓迈   ·   智能座舱创新 DEMO 样机解决方案',EN,11,True,CYAN)]])
txt(s,0.7,1.7,5.4,1.5,[[('把硬件优势,变成',CN,31,True,WHITE)],[('主机厂面前的',CN,31,True,WHITE)],[('体验话语权',CN,31,True,LILAC)]],ls=1.05)
pic(s,5.6,1.55,3.9,2.7,'/tmp/assets/cover.jpg',border=RGBColor(0x35,0x46,0x80))
txt(s,5.6,4.28,3.9,0.3,[[('卓迈 iX-2024 · 红点奖座舱',CN,8.5,True,CYAN)]],al=CT)
txt(s,0.72,3.7,5.0,0.4,[[('NICE 智能座舱体验舱 · 整包样机方案',CN,13,False,FOG)]])
rect(s,0.72,4.2,2.6,0.025,fill=RGBColor(0x35,0x46,0x80))
txt(s,0.7,4.55,5.0,0.4,[[('呈送:宁波华翔',CN,12,True,WHITE),('   面向 Tier 1',CN,11,False,MUTE)]])
txt(s,0.7,5.02,8.8,0.3,[[('上海卓迈 JOMEC   |   2026.06   |   机密·仅限指定收件方',CN,9.5,False,MUTE)]])
pg()

# ============ 2 目录
s,_=slide(WHITE); chead(s,pg(),'目录:六段旅程,一条完整说服链','CONTENTS · THE WINNING PATH')
parts=[('01','你的处境与机会','主机厂要求抬升 · L1→L4 能力阶梯','target'),('02','我们是谁','一站式责任主体 · 德系×中国速度','hand'),
       ('03','样机解决方案','样机真假 · 量产级工程底座 · 专利库','cube'),('04','落到 NICE 体验舱','4 场景 · 系统架构 · 30 周计划 · 商务','car'),
       ('05','验证过的赢法','案例证据链 · 红点/宝马/Apollo','trophy'),('06','你将得到什么','12 个月后状态 · 投资逻辑 · 下一步','check')]
for i,(n,t,d,ic) in enumerate(parts):
    x=C2[i%2]; y=1.2+(i//2)*1.2; ac=[BLUE,CYAN,MAG,BLUE,CYAN,VIO][i]
    ccard(s,x,y,W2,1.02,topbar=ac)
    circ(s,x+0.6,y+0.51,0.66,ac,fill2=[VIO,BLUE,VIO,CYAN,BLUE,MAG][i],icon=ic)
    txt(s,x+1.08,y+0.16,2.6,0.4,[[(n+'  ',EN,12,True,ac),(t,CN,14,True,TITLED)]])
    txt(s,x+1.08,y+0.6,3.2,0.35,[[(d,CN,8.5,False,BODY)]])
cfoot(s,P[0])

# ============ 3 PART01
divider('PART 01','你的处境与机会','THE NEW STANDARD FROM OEMs',
 '主机厂正把『用户体验创新』写进定点评判标准,并向内外饰一级供应商传导。补齐『前瞻造型—交互—电子软硬件—可落地样机』这条能力链,已不是加分项,而是守住订单的前提。','01'); pg()

# ============ 4 L1-L4
s,_=slide(WHITE); chead(s,pg(),'行业拐点:主机厂要求,从『供零件』走向『供体验』','THE BAR IS RISING')
story(s,'老轮机长贴着机舱壁一听就知道哪个缸不对——主机厂的『耳朵』也在变:从听质量,到听你懂不懂体验。',label='趋势',accent=BLUE)
levels=[('L1','供零件','按图交付合格零部件','成本·质量·交期',BLUE),('L2','供模块','交付集成化总成模块','集成度·可靠性·降本',CYAN),
        ('L3','供体验','参与定义座舱体验创新','前瞻造型·交互·电子·机构',MAG),('L4','供创新合作','成为主机厂共创伙伴','主导话题·带样机进评审',VIO)]
for i,(l,t,d,k,c) in enumerate(levels):
    x=C4[i]; y=1.88; hl=(l in('L3','L4'))
    cc=ccard(s,x,y,W4,2.32,topbar=c)
    if hl: rect(s,x,y,W4,2.32,line=MAG,lw=1.4,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
    txt(s,x+0.2,y+0.16,W4-0.4,0.5,[[(l,ENB,22,True,c)]])
    txt(s,x+0.2,y+0.7,W4-0.4,0.35,[[(t,CN,14,True,TITLED)]])
    txt(s,x+0.2,y+1.1,W4-0.4,0.7,[[(d,CN,9,False,BODY)]],ls=1.15)
    txt(s,x+0.2,y+1.86,W4-0.4,0.4,[[(k,CN,8,True,c)]],ls=1.1)
bottom(s,'华翔今天正卡在 L2 → L3 的台阶上','制造与模块能力扎实,但『供体验』所需的综合能力,正是缺口。')
cfoot(s,P[0])

# ============ 5 缺口补位
s,_=slide(WHITE); chead(s,pg(),'迈向『供体验』:你要补齐的四块 + 一个保障','THE GAP & HOW JOMEC FILLS IT')
rows=[('前瞻造型','德系履历造型团队 + CAS 快速建模','拿得出让主机厂眼前一亮的造型语言','palette'),
      ('交互 / UI / UX','完整交互系统开发(UI/UX/动效/VR)','正中主机厂体验创新的新需求','touch'),
      ('电子软硬件','自研车规级控制器 + 嵌入式集成','软硬件自主可控,不被第三方卡脖子','chip'),
      ('运动机构','量产级机构(折叠方向盘达上汽量产)','最难的机构有人兜底,还能上车','gears'),
      ('责任归属(保障)','五位一体单一责任主体','你只对接一家,不背协同的锅','hand')]
txt(s,1.35,1.12,2.4,0.3,[[('要补的能力',CN,9.5,True,MUTE)]])
txt(s,3.7,1.12,3.0,0.3,[[('JOMEC 补齐',CN,9.5,True,BLUE)]])
txt(s,6.8,1.12,2.7,0.3,[[('对你意味着',CN,9.5,True,MAG)]])
for i,(a,b,c,ic) in enumerate(rows):
    y=1.45+i*0.585; hl=(i==4); ac=MAG if hl else BLUE
    cc=ccard(s,0.55,y,8.9,0.52,topbar=ac)
    circ(s,0.92,y+0.26,0.42,ac,fill2=VIO if hl else CYAN,icon=ic)
    txt(s,1.3,y,2.4,0.52,[[(a,CN,11,True,TITLED)]],an=MID)
    txt(s,3.7,y,3.0,0.52,[[(b,CN,9.5,False,BODY)]],an=MID)
    txt(s,6.8,y,2.6,0.52,[[(c,CN,9.5,True,(MAG if hl else BODY))]],an=MID)
cfoot(s,P[0])

# ============ 6 PART02
divider('PART 02','我们是谁','COMPANY & CAPABILITIES',
 '卓迈 JOMEC,国内少数能同时交付造型、交互、电子与量产级机构的创新公司。成立于 2011 年,核心团队来自国内主流车企与科技公司,既懂德系高标准,又懂中国量产节奏。','02'); pg()

# ============ 7 公司概况
s,_=slide(WHITE); chead(s,pg(),'公司概况:造型 + 交互 + 电子 + 量产级机构','COMPANY PROFILE')
bases=[('上海','前瞻中心'),('昆山','创新研发&试制'),('柳州','数字化中心'),('杭州','JOMEC×CAA'),('宁波','产品生产基地')]
txt(s,0.55,1.08,8.9,0.3,[[('五大基地 · 一条龙能力',CN,10,True,BLUE)]])
for i,(c,d) in enumerate(bases):
    x=0.55+i*1.78; bw=1.7; ac=[BLUE,CYAN,VIO,MAG,BLUE][i]
    ccard(s,x,1.42,bw,1.28,topbar=ac)
    circ(s,x+bw/2,1.86,0.5,ac,fill2=VIO,icon=None,label='',)
    lp=AS%LOC[i]
    if os.path.exists(lp):
        s.shapes.add_picture(lp,Inches(x+bw/2-0.13),Inches(1.73),Inches(0.26),Inches(0.26))
    txt(s,x,2.18,bw,0.3,[[(c,CN,13,True,TITLED)]],al=CT)
    txt(s,x+0.08,2.46,bw-0.16,0.3,[[(d,CN,7.5,False,BODY)]],al=CT,ls=1.05)
ccard(s,0.55,2.92,4.35,1.32,topbar=BLUE)
txt(s,0.8,3.08,3.9,0.35,[[('开发业务',CN,11,True,BLUE)]])
txt(s,0.8,3.45,3.9,0.7,[[('未来出行造型开发 · 交互设计 · 智能座舱创新开发',CN,10,False,BODY)]],ls=1.25)
ccard(s,5.1,2.92,4.35,1.32,topbar=CYAN)
txt(s,5.35,3.08,3.9,0.35,[[('产品业务',CN,11,True,CYAN)]])
txt(s,5.35,3.45,3.9,0.7,[[('座舱智能总成产品研发及供货 · 量产交付',CN,10,False,BODY)]],ls=1.25)
bottom(s,'卓迈使命:塑造更美好的未来生活','始终创业 · 求真务实 · 敢为极致 · 共同成长')
cfoot(s,P[0])

# ============ 8 服务全景
s,_=slide(WHITE); chead(s,pg(),'服务全景:从调研到生产,覆盖创新开发全链条','OUR SERVICES')
steps=[('01','交通工具调研','竞品·趋势·用户·策略','search'),('02','交通工具设计','概念·内外饰·CMF·A级曲面','palette'),
       ('03','展车与原型','CNC·3D打印·油泥·功能样机','cube'),('04','UX-UI 设计','HMI·体验·动效·软件·VR/AR','touch'),('05','生产制造','快速原型·功能样件·量产','factory')]
for i,(n,t,d,ic) in enumerate(steps):
    x=0.55+i*1.86; y=1.42; ac=[BLUE,CYAN,VIO,MAG,BLUE][i]
    ccard(s,x,y,1.72,2.4,topbar=ac)
    circ(s,x+0.86,y+0.6,0.66,ac,fill2=VIO,icon=ic)
    txt(s,x+0.06,y+1.05,1.6,0.4,[[(n+' ',EN,10,True,ac),(t,CN,10.5,True,TITLED)]],al=CT,ls=1.05)
    txt(s,x+0.06,y+1.6,1.6,0.7,[[(d,CN,8,False,BODY)]],al=CT,ls=1.15)
    if i<4: txt(s,x+1.7,y+0.5,0.2,0.5,[[('▸',EN,12,True,MUTE)]])
bottom(s,'调研 → 设计 → 原型 → 体验 → 制造','创新不会断在任何一个交接环节。')
cfoot(s,P[0])

# ============ 9 团队
s,_=slide(WHITE); chead(s,pg(),'团队:德系深度 × 中国速度','TEAM')
ppl=[('严 艇 Tim','CEO · 中国美院客座教授','江苏省双创人才;AI 驱动座舱创新方向'),
     ('黄 斌 Burt','联合创始人 · 造型','前泛亚汽车技术中心设计执行总监;25 年经验'),
     ('Alexander Zeth','联合创始人 · 设计','德籍;前 EDAG 资深设计师;服务宝马/奥迪/大众'),
     ('徐康聪','总工程师','德国布伦瑞克博士后;两获中国汽车工业科学进步奖')]
for i,(n,r,d) in enumerate(ppl):
    x=C2[i%2]; y=1.15+(i//2)*1.05; ac=[BLUE,CYAN,VIO,MAG][i]
    ccard(s,x,y,W2,0.92,topbar=ac)
    circ(s,x+0.5,y+0.5,0.56,ac,fill2=VIO,label=n[0],lsize=15)
    txt(s,x+0.88,y+0.13,3.4,0.32,[[(n,CN,12,True,TITLED)]])
    txt(s,x+0.88,y+0.43,3.4,0.28,[[(r,CN,8.5,True,BLUE)]])
    txt(s,x+0.88,y+0.68,3.45,0.24,[[(d,CN,7.5,False,BODY)]])
chips=[('IATF16949 全量产流程',BLUE),('中国美院 CAA 联合前瞻',CYAN),('两度赴德:宝马/奔驰/大众/奥迪',MAG)]
for i,(t,c) in enumerate(chips):
    x=0.55+i*2.98; rr=rect(s,x,3.4,2.85,0.42,fill=PANEL,line=c,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5); shadow(rr,blur=40000,dist=12000,alpha=85000)
    txt(s,x,3.4,2.85,0.42,[[(t,CN,8.5,True,c)]],al=CT,an=MID)
bottom(s,'既懂德系高标准,又懂中国量产节奏','这是大多数本土与外资团队都给不齐的组合。')
cfoot(s,P[0])

# ============ 10 组织(放射图·满铺暗底 hero)
s,bgr=slide(DARK)
n10=pg()
if os.path.exists('/tmp/assets/hub.png'):
    s.shapes.add_picture('/tmp/assets/hub.png',0,0,prs.slide_width,prs.slide_height)
txt(s,8.55,5.31,0.85,0.24,[[(f'{n10:02d}',EN,8,True,RGBColor(0x8A,0x98,0xB6))]],al=RT)

# ============ 11 PART03
divider('PART 03','样机解决方案','THE DEMO SOLUTION',
 '样机有真假之分,主机厂评审现场见分晓。能摸、能动、能现场联调、过信号干涉测试的 L3 全功能可交互样机,才扛得住工程拷问。','03'); pg()

# ============ 12 样机真假
s,_=slide(WHITE); chead(s,pg(),'重新定义品类:样机有真假之分','REAL vs FAKE DEMO')
story(s,'好看的壳子 + 假屏,评委一问『能动吗』就卡住;唯有能摸能动、过信号干涉测试的样机,才取信工程评委。',label='评审现场',accent=MAG)
lv=[('L1','静态造型舱','好看的壳子 + 假屏。只能看、能拍照;一问能动吗就卡住。',MUTE,RGBColor(0xC2,0xCB,0xDB)),
    ('L2','数字孪生 / 视频','屏幕里的炫酷演示。只能播放;摸不到、难取信评委。',MUTE,RGBColor(0xC2,0xCB,0xDB)),
    ('L3','全功能可交互样机','能摸、能动、能现场联调、过信号干涉测试。扛工程拷问。',MAG,MAG)]
for i,(l,t,d,tc,bc) in enumerate(lv):
    x=C3[i]; hl=(i==2)
    ccard(s,x,1.9,W3,1.78,topbar=bc)
    if hl: rect(s,x,1.9,W3,1.78,line=MAG,lw=1.4,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
    txt(s,x+0.2,2.06,W3-0.4,0.45,[[(l,ENB,18,True,tc)]])
    txt(s,x+0.2,2.55,W3-0.4,0.35,[[(t,CN,12,True,TITLED)]])
    txt(s,x+0.2,2.95,W3-0.4,0.7,[[(d,CN,8.5,False,BODY)]],ls=1.2)
bottom(s,'JOMEC 的定义:能力闭环 · 单一责任主体','造型→交互→软硬件→机构→样机,协同出的问题我们兜底。')
cfoot(s,P[0])

# ============ 13 工程底座
s,_=slide(WHITE); chead(s,pg(),'差异化证据:能动样机背后,是量产级工程','THE ENGINEERING BASE')
eng=[('测试台架','把风险提前到样机阶段','样机进评审室前先在台架把问题出完:总线/信号干涉/氛围灯/机构/多屏/散热/识别','9 项联调验证',BLUE,'/tmp/assets/bench.jpg'),
     ('运动机构','最难的开闭件已到量产级','折叠方向盘达上汽乘用车量产级并装车联调;鸥翼门/侧门应用于百度 Apollo','上汽量产级',CYAN,'/tmp/assets/seat.jpg'),
     ('自研车规控制器','软硬件是自己的,不外包','已自研 5 类实物控制器:通用/多路通断/RGB氛围灯/三相电机/小型舵机','5 类自研控制器',MAG,'/tmp/assets/ctrl.jpg')]
for i,(t,sub,d,tag,c,im) in enumerate(eng):
    x=C3[i]; ccard(s,x,1.2,W3,3.0,topbar=c)
    pic(s,x+0.15,1.33,W3-0.3,1.0,im,border=c)
    txt(s,x+0.2,2.44,W3-0.4,0.32,[[(t,CN,12,True,TITLED),('  '+sub,CN,8.5,True,c)]])
    txt(s,x+0.2,2.82,W3-0.4,0.95,[[(d,CN,8.5,False,BODY)]],ls=1.22)
    tg=rect(s,x+0.2,3.74,W3-0.4,0.38,fill=PANEL,line=c,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.4)
    txt(s,x+0.2,3.74,W3-0.4,0.38,[[(tag,CN,9,True,c)]],al=CT,an=MID)
bottom(s,'反复出现的『量产级』『车规』『通过性能测试』','这是 JOMEC 与造型模型公司的根本分野。')
cfoot(s,P[0])

# ============ 14 专利库
s,_=slide(WHITE); chead(s,pg(),'降低决策风险:自有专利创新库,即插即用','OWNED-PATENT LIBRARY')
pat=[('透光表面 ShyTech','去物理按键仍保留全部开关','900–1200lm/m²·寿命5万h·透光92%',MAG,'华翔主场','bulb'),
     ('隐藏式电动长滑轨','把『第三空间』做出来','承重15kg·收纳740×560×270·30°爬坡',BLUE,'','seat'),
     ('椅背磁吸智联拓展','屏幕/平板即吸即用+无线充','90°旋转·磁吸+无线充·减少凸露',CYAN,'','magnet'),
     ('重力感应升降杯托','自动感知适应物品重量','寿命≥10万次·IPX4·20°不洒',VIO,'','cup'),
     ('折叠方向盘机构','释放驾驶位空间','已达上汽量产级·与线控转向同步',BLUE,'','gears')]
for i,(t,sub,spec,c,tag,ic) in enumerate(pat):
    if i<3: x=C3[i]; y=1.15
    else: x=[2.04,5.51][i-3]; y=2.75
    ccard(s,x,y,W3,1.45,topbar=c)
    circ(s,x+0.42,y+0.42,0.5,c,fill2=VIO,icon=ic)
    txt(s,x+0.78,y+0.13,(W3-1.7 if tag else W3-0.9),0.55,[[(t,CN,10.5,True,TITLED)]],an=MID)
    if tag:
        rect(s,x+W3-0.95,y+0.16,0.8,0.28,fill=NAVY,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
        txt(s,x+W3-0.95,y+0.16,0.8,0.28,[[(tag,CN,7.5,True,LILAC)]],al=CT,an=MID)
    txt(s,x+0.18,y+0.7,W3-0.3,0.3,[[(sub,CN,8.5,True,c)]])
    txt(s,x+0.18,y+1.02,W3-0.3,0.4,[[(spec,CN,7.5,False,BODY)]],ls=1.1)
bottom(s,'这些是 JOMEC 自有专利','你拿到的是干净、可申报的创新点,而不是侵权风险。',y=4.34,h=0.5)
cfoot(s,P[0])

# ============ 15 交付路径
s,_=slide(WHITE); chead(s,pg(),'交付路径:从概念到可演示样机,每阶段都验收到弹药','DELIVERY PATH')
ph=[('1','概念阶段','创新切入点定义 · CAS 造型 · 交互原型','可沟通:方向与亮点',BLUE),
    ('2','工程阶段','运动机构方案 · 线束 3D · 软硬件架构','可沟通:可行性与深度',CYAN),
    ('3','样机阶段','台架联调通过 · 全功能可交互样机','可沟通:能现场演示的实物',MAG)]
for i,(n,t,d,say,c) in enumerate(ph):
    x=C3[i]; ccard(s,x,1.25,W3,2.85,topbar=c)
    circ(s,x+0.55,1.78,0.66,c,fill2=VIO,label=n,lsize=18)
    txt(s,x+1.0,1.55,W3-1.1,0.45,[[(t,CN,13,True,TITLED)]],an=MID)
    txt(s,x+0.2,2.4,W3-0.4,0.9,[[(d,CN,9,False,BODY)]],ls=1.25)
    rr=rect(s,x+0.2,3.42,W3-0.4,0.55,fill=PANEL,line=c,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.15)
    txt(s,x+0.3,3.42,W3-0.6,0.55,[[(say,CN,8.5,True,c)]],an=MID,ls=1.1)
    if i<2: txt(s,x+W3-0.05,2.05,0.3,0.5,[[('▸',EN,12,True,MUTE)]])
bottom(s,'UX 研究流程与 IATF16949 工程流程融合为一条主线','每个节点都翻译成你能对主机厂说的话。')
cfoot(s,P[0])

# ============ 16 PART04
divider('PART 04','落到 NICE:体验舱整包方案','THE NICE EXPERIENCE CABIN',
 'NICE 智能座舱体验舱,就是 L3 全功能可交互样机能力的旗舰落地。以下是我们对华翔这个具体项目已经想到的深度 —— 从体验场景、系统架构、安全可靠,到 30 周交付与商务。','04'); pg()

# ============ 17 NICE概念
s,_=slide(WHITE); chead(s,pg(),'NICE 体验舱:健康舒适 + 智能表面','THE EXPERIENCE CABIN')
ccard(s,0.55,1.2,5.0,3.0,topbar=BLUE)
txt(s,0.78,1.4,4.6,0.4,[[('主题:',CN,11,True,BLUE),('健康舒适 + 智能表面',CN,14,True,TITLED)]])
fe=['可进入 — 支持 4 人同时试乘','可试乘 — 4 套零重力座椅(前排含 180° 旋转)','可交互 — 语音 + 多屏触控 + 实体按键',
    '可循环 — 日连续运行 ≥10 小时无 A 类故障','可迭代 — 预留扩展,支撑全球巡展']
for i,f in enumerate(fe):
    y=1.95+i*0.44; rect(s,0.78,y+0.05,0.12,0.12,fill=[BLUE,CYAN,LILAC,BLUE,CYAN][i])
    txt(s,1.02,y,4.4,0.4,[[(f,CN,9.5,False,BODY)]])
if os.path.exists('/tmp/assets/ix.jpg'):
    ph=s.shapes.add_picture('/tmp/assets/ix.jpg',Inches(5.75),Inches(1.2),Inches(3.7),Inches(2.05)); shadow(ph,blur=55000,dist=18000,alpha=78000)
    txt(s,5.75,3.32,3.7,0.32,[[('卓迈 iX-2024 红点奖座舱',CN,9,True,BLUE)]],al=CT)
    txt(s,5.75,3.66,3.7,0.5,[[('同类获奖座舱示意 · 待替换为 NICE 专属渲染',CN,8,False,MUTE)]],al=CT,ls=1.1)
else:
    ph=rect(s,5.75,1.2,3.7,3.0,fill=PANEL,line=CLINE,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05); shadow(ph)
    txt(s,5.75,2.85,3.7,0.4,[[('〔 体验舱效果图占位 〕',CN,11,True,MUTE)]],al=CT)
bottom(s,'1:1 可进入 · 可试乘 · 可交互 · 可循环演示 · 可持续迭代')
cfoot(s,P[0])

# ============ 18 四场景
s,_=slide(WHITE); chead(s,pg(),'四大体验:一次进入,四重惊艳','FOUR WOW MOMENTS')
sc=[('S1','迎宾','走近即点亮——星河迎宾+氛围灯流水+星空天幕','P0',BLUE,'bulb'),('S2','健康舒适','坐下被照顾——零重力+健康可视化+全车加热','P0',CYAN,'eye'),
    ('S3','智能表面','无屏即屏——IMSE膜+自发光缝线+材质质感','P1',MAG,'finger'),('S4','智能空间','一键变形——Console展开+座椅180°旋转','P1',VIO,'seat')]
for i,(n,t,d,p,c,ic) in enumerate(sc):
    x=C4[i]; ccard(s,x,1.2,W4,2.95,topbar=c)
    circ(s,x+0.46,1.66,0.56,c,fill2=VIO,icon=ic)
    txt(s,x+0.84,1.44,W4-0.9,0.45,[[(n+' ',ENB,12,True,c),(t,CN,12.5,True,TITLED)]],an=MID)
    pr=rect(s,x+0.18,2.12,0.55,0.28,fill=PANEL,line=(RGBColor(0xF4,0x3F,0x5E) if p=='P0' else CYAN),sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,x+0.18,2.12,0.55,0.28,[[(p,EN,8,True,(RGBColor(0xF4,0x3F,0x5E) if p=='P0' else CYAN))]],al=CT,an=MID)
    txt(s,x+0.18,2.55,W4-0.34,1.5,[[(d,CN,8.5,False,BODY)]],ls=1.3)
bottom(s,'每个场景都是触发→系统响应→技术亮点→复位完整闭环','覆盖 ≥2 种交互方式。')
cfoot(s,P[0])

# ============ 19 架构安全(分层图)
s,_=slide(WHITE); chead(s,pg(),'系统架构 & 安全可靠:量产级方法治理复杂度','ARCHITECTURE & SAFETY')
layers=[('演示/运维层','演示平板 · 后台调试 · 硬件复位',BLUE),
        ('中央控制层','域控制器(主控) · 多屏 + 5×CAN FD 网关',VIO),
        ('显示/视频层','P-HUD · 魔术屏 · 吸顶屏 · 后排屏 · PAD',CYAN),
        ('域总线层 CAN FD','灯光 / 座椅 / 机构 / 舒适 / 加热健康 · ≥5Mbps',BLUE),
        ('节点执行层','控制板 · 驱动器 · 加热膜 · 健康扶手',VIO),
        ('配电安全层','PSU×5 · 急停 · 限位 · 防夹 · 超温自断',MAG)]
for i,(t,d,c) in enumerate(layers):
    y=1.15+i*0.5; rr=rect(s,0.55,y,5.9,0.44,fill=PANEL,line=CLINE,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.18); shadow(rr,blur=35000,dist=10000,alpha=88000)
    rect(s,0.55,y,0.9,0.44,fill=c,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.18); rect(s,1.05,y,0.4,0.44,fill=c)
    txt(s,0.55,y,0.9,0.44,[[('L'+str(6-i),ENB,12,True,WHITE)]],al=CT,an=MID)
    txt(s,1.6,y,2.0,0.44,[[(t,CN,9.5,True,TITLED)]],an=MID)
    txt(s,3.55,y,2.8,0.44,[[(d,CN,8,False,BODY)]],an=MID)
# 右侧安全要点卡
saf=[('硬件安全回路','急停+STO+限位+电流防夹,独立于软件',MAG),('稳定与兜底','≥10h连续·100次无故障·一键复位≤30s',CYAN),('断电恢复','上电自动回零+参数记忆,≤5分钟',BLUE)]
for i,(t,d,c) in enumerate(saf):
    y=1.15+i*1.0; ccard(s,6.65,y,2.8,0.88,topbar=c)
    txt(s,6.85,y+0.13,2.5,0.35,[[(t,CN,10.5,True,TITLED)]])
    txt(s,6.85,y+0.46,2.5,0.4,[[(d,CN,7.5,False,BODY)]],ls=1.12)
bottom(s,'真实电气工程,而非模型搭建','这套架构让"展会不出事、A 类零故障"的承诺有出处。')
cfoot(s,P[0])

# ============ 20 30周(甘特+网格)
s,_=slide(WHITE); chead(s,pg(),'30 周交付确定性:2026.07 启动 → 2027.01 终验收','30-WEEK ROADMAP')
gx,gy,gw=2.0,1.75,7.2
mons=[('07',2),('08',2),('09',2),('10',2),('11',2),('12',2),('27.01',1.5)]
tot=sum(m[1] for m in mons); ux=gw/tot; cx=gx
gh=0.4+6*0.4
for m,wd in mons:
    rect(s,cx,gy-0.3,ux*wd,0.26,fill=PANEL,line=CLINE); txt(s,cx,gy-0.3,ux*wd,0.26,[[(m,EN,8,True,TITLED)]],al=CT,an=MID)
    rect(s,cx,gy,0.008,2.45,fill=RGBColor(0xEC,0xEF,0xF6))  # 网格竖线
    cx+=ux*wd
teams=[('造型/CAS',0.0,0.52,BLUE),('硬件开发',0.08,0.74,RGBColor(0xF4,0x3F,0x5E)),('软件开发',0.16,0.80,CYAN),
       ('展车模型',0.20,0.80,RGBColor(0xF5,0x9E,0x0B)),('动画/交互',0.0,0.72,VIO),('测试(前置)',0.60,0.40,RGBColor(0x22,0xC5,0x5E))]
for i,(t,st,ln,c) in enumerate(teams):
    y=gy+0.12+i*0.4
    txt(s,0.55,y,1.4,0.32,[[(t,CN,9,True,TITLED)]],an=MID)
    bar=rect(s,gx+st*gw,y+0.05,ln*gw,0.22,fill=c,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.4); shadow(bar,blur=30000,dist=8000,alpha=85000)
ms=[('M4 骨架',0.30),('M6 预验收',0.66),('M7 终验收',0.93)]
for t,fr in ms:
    x=gx+fr*gw; rect(s,x-0.08,gy+0.0,0.16,0.16,fill=MAG,sh=MSO_SHAPE.DIAMOND)
    txt(s,x-0.6,4.28,1.2,0.3,[[(t,CN,8,True,MAG)]],al=CT)
bottom(s,'并行工程 + 测试 W18 前置','◆ 为里程碑(执行计划口径,合同关键节点以正式合同为准)。')
cfoot(s,P[0])

# ============ 21 商务
s,_=slide(WHITE); chead(s,pg(),'商务与服务范围:价值而非成本','COMMERCIAL & SERVICE')
svc=['免费质保期 ≥ 1 年','展会驻场 ≥2 人 · A 类响应 ≤15min','远程支持 7×24 · ≤2 小时','免费软件迭代 ≤3 次','关键易损件供应 ≥3 年','质保期展会支持 ≥10 次/16 天']
ccard(s,0.55,1.2,4.5,3.0,topbar=BLUE)
txt(s,0.78,1.38,4.1,0.35,[[('服务承诺',CN,11,True,BLUE)]])
for i,f in enumerate(svc):
    y=1.78+i*0.4; rect(s,0.78,y+0.05,0.12,0.12,fill=[BLUE,CYAN,LILAC][i%3]); txt(s,1.0,y,3.9,0.4,[[(f,CN,9.5,False,BODY)]])
ccard(s,5.25,1.2,4.2,3.0,topbar=MAG)
txt(s,5.5,1.38,3.8,0.35,[[('报价结构(占位)',CN,11,True,MAG)]])
quote=['造型与 CMF 设计','工程与展车制作','电子电器与机构','软件/HMI/集成','动画与传播','展会保障与售后']
for i,q in enumerate(quote):
    y=1.82+i*0.38; txt(s,5.5,y,2.6,0.33,[[(q,CN,9,False,BODY)]])
    rect(s,8.25,y+0.02,1.0,0.27,fill=PANEL,line=CLINE); txt(s,8.25,y+0.02,1.0,0.27,[[('〔待报〕',CN,8,False,MUTE)]],al=CT,an=MID)
bottom(s,'可按概念/工程/样机分阶段推进','每阶段有可验收交付物;具体投入基于项目范围当面测算。')
cfoot(s,P[0])

# ============ 22 PART05
divider('PART 05','验证过的赢法','PROVEN CASES',
 '这套打法,真把别的 Tier 1 送进了主机厂的门。按说服力排序:从带着创新走进宝马,到自有专利获捷豹路虎官方致谢,再到红点奖与百度 Apollo —— 同一套能力可跨客户、跨场景复用。','05'); pg()

# ============ 23 案例证据链
s,_=slide(WHITE); chead(s,pg(),'案例证据链:按说服力排序','EVIDENCE, RANKED')
cs=[('常熟汽饰 × 宝马','主导 iX-2024 隐藏式电动长滑轨创新,负责对接与提案','宝马创新开发契机 + 量产机会',MAG,'trophy'),
    ('椅背磁吸 × 捷豹路虎','凭 JOMEC 自有专利总成做技术展示','JLR 认可 + 官方感谢信',BLUE,'magnet'),
    ('2024 德国红点奖','JOMEC 申报;iX-2024 两度赴德独立展示','行业级设计认可',CYAN,'check'),
    ('百度 Apollo 无人车','底盘/造型/工程/硬件/HMI/样车制作','入藏北京汽车博物馆',VIO,'car')]
for i,(t,d,r,c,ic) in enumerate(cs):
    x=C2[i%2]; y=1.2+(i//2)*1.5; ccard(s,x,y,W2,1.32,topbar=c)
    circ(s,x+0.45,y+0.45,0.52,c,fill2=VIO,icon=ic)
    txt(s,x+0.82,y+0.14,W2-0.9,0.4,[[(t,CN,12,True,TITLED)]])
    txt(s,x+0.82,y+0.54,W2-0.95,0.45,[[(d,CN,8.5,False,BODY)]],ls=1.1)
    txt(s,x+0.2,y+0.98,W2-0.4,0.3,[[('▸ ',EN,9,True,c),(r,CN,9,True,c)]])
bottom(s,'「 我们让一家 Tier 1,带着创新走进了宝马的门。 」','这是 PPT 讲不出来的说服力。')
cfoot(s,P[0])

# ============ 24 重点案例(船舶实拍)
s,_=slide(WHITE); chead(s,pg(),'重点案例:JoySpace+ 与 船舶智慧驾驶舱','FLAGSHIP CASES')
ccard(s,0.55,1.2,4.4,3.0,topbar=BLUE)
if os.path.exists('/tmp/assets/joy.jpg'):
    jp=s.shapes.add_picture('/tmp/assets/joy.jpg',Inches(5.2 if False else 0.7),Inches(1.32),Inches(4.1),Inches(1.45)); shadow(jp,alpha=80000)
txt(s,0.8,2.86,4.0,0.32,[[('JoySpace+ × 均胜电子',CN,12,True,TITLED)]])
for i,f in enumerate(['沉浸式智能座舱·头部Tier1共创范式','赛博机能美学·天空之镜水晶极光·AI智慧体','折叠方向盘·电子后视镜·星空顶篷']):
    y=3.26+i*0.32; rect(s,0.8,y+0.04,0.1,0.1,fill=CYAN); txt(s,1.0,y,3.85,0.32,[[(f,CN,8,False,BODY)]])
ccard(s,5.05,1.2,4.4,3.0,topbar=CYAN)
if os.path.exists('/tmp/assets/ship_crop.jpg'):
    pic=s.shapes.add_picture('/tmp/assets/ship_crop.jpg',Inches(5.2),Inches(1.32),Inches(4.1),Inches(1.5)); shadow(pic,alpha=80000)
txt(s,5.3,2.92,4.0,0.35,[[('船舶智慧驾驶舱',CN,12,True,TITLED)]])
for i,f in enumerate(['七屏中枢+数字孪生·180°旋转座椅','人脸/手势/无感健康监测/AI语音','滑动方向舵·智慧海图桌·自适应氛围灯']):
    y=3.32+i*0.32; rect(s,5.3,y+0.04,0.1,0.1,fill=LILAC); txt(s,5.5,y,3.85,0.32,[[(f,CN,8,False,BODY)]])
bottom(s,'从汽车到船舶,同一套造型+交互+软硬件+样机能力','证明创新方法论可迁移。')
cfoot(s,P[0])

# ============ 25 生态
s,_=slide(WHITE); chead(s,pg(),'战略合作伙伴生态,提升项目技术含量','PARTNER ECOSYSTEM')
pt=[('丰田通商','世界 500 强','新技术新材料及开发支持',BLUE),('CP-VELA','国际气联认证工程服务','新技术新材料开发支持',CYAN),
    ('百度 Apollo','生态合作(燎原计划)','已达成生态合作伙伴关系',VIO),('科思创 Covestro','全球领先聚合物供应商','聚氨酯蜂窝/红外透波PC材料',MAG)]
for i,(t,r,d,c) in enumerate(pt):
    x=C4[i]; ccard(s,x,1.4,W4,2.5,topbar=c)
    circ(s,x+W4/2,1.98,0.66,c,fill2=VIO,label=t[0],lsize=16)
    txt(s,x+0.1,2.5,W4-0.2,0.4,[[(t,CN,11,True,TITLED)]],al=CT)
    txt(s,x+0.1,2.9,W4-0.2,0.3,[[(r,CN,8,True,c)]],al=CT)
    txt(s,x+0.12,3.22,W4-0.24,0.55,[[(d,CN,7.5,False,BODY)]],al=CT,ls=1.15)
bottom(s,'生态合作不是背书装点','而是把新材料、新技术真正带进华翔的项目。')
cfoot(s,P[0])

# ============ 26 PART06
divider('PART 06','你将得到什么','YOUR FUTURE & NEXT STEP',
 '12 个月后,华翔将不再是被动接收主机厂创新需求的硬件供应商,而是带着一台经得起拷问的样机、主导评审现场体验话题的创新合作方。','06'); pg()

# ============ 27 12个月后
s,_=slide(WHITE); chead(s,pg(),'12 个月后的宁波华翔','YOUR FUTURE STATE')
ccard(s,0.55,1.2,4.4,3.0,topbar=BLUE)
txt(s,0.8,1.38,4.0,0.4,[[('你将拿到的(可验收)',CN,12,True,TITLED)]])
for i,f in enumerate(['1 台全功能可交互样机','造型数据 + 交互原型','软硬件方案 + 运动机构','台架联调报告','可对外演示脚本']):
    y=1.85+i*0.45; rect(s,0.8,y+0.05,0.12,0.12,fill=BLUE); txt(s,1.02,y,3.8,0.4,[[(f,CN,10,False,BODY)]])
cc=rect(s,5.05,1.2,4.4,3.0,fill=PANEL,line=CLINE,sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05); shadow(cc,blur=50000,dist=16000,alpha=82000); rect(s,5.05,1.2,0.05,3.0,fill=MAG)
txt(s,5.3,1.38,4.0,0.4,[[('竞争位置的转变',CN,12,True,MAG)]])
txt(s,5.3,1.95,4.0,0.4,[[('从',CN,10,True,MUTE),(' 被动接收体验需求的硬件供应商',CN,10,False,BODY)]],ls=1.2)
txt(s,5.3,2.65,4.0,0.3,[[('↓',EN,14,True,MAG)]])
txt(s,5.3,3.05,4.0,0.7,[[('到',CN,10,True,BLUE),(' 带着能动样机、主导评审现场体验话题的创新合作方',CN,10,True,TITLED)]],ls=1.25)
bottom(s,'下一次主机厂问"你们能做体验创新吗"','华翔的回答不是 PPT,是一台能让评委坐进去的样机。')
cfoot(s,P[0])

# ============ 28 投资逻辑
s,_=slide(WHITE); chead(s,pg(),'投资逻辑:一次定点的价值,远大于一台样机的投入','VALUE, NOT COST')
inv=[('你买的不是设计工时','而是 L3 全功能可交互样机 —— 不是 L1 造型舱,也不是 L2 视频',BLUE),
     ('你买的不是一次性概念','而是自有专利创新库 —— IP 干净、可复用、规避侵权风险',CYAN),
     ('你买的不是本土概念能力','而是德系量产级标准 + 中国速度',MAG)]
for i,(a,b,c) in enumerate(inv):
    y=1.3+i*1.0; ccard(s,0.55,y,8.9,0.85,topbar=c)
    txt(s,0.85,y,3.3,0.85,[[(a,CN,12,True,TITLED)]],an=MID)
    txt(s,4.3,y,5.0,0.85,[[(b,CN,10.5,False,BODY)]],an=MID,ls=1.2)
bottom(s,'用一台样机的成本,守住或赢得一次主机厂定点的机会','把投入放回它真正的语境。')
cfoot(s,P[0])

# ============ 29 下一步(暗·呼应封面)
s,bgr=slide(NAVY); grad(bgr,RGBColor(0x07,0x14,0x33),RGBColor(0x1A,0x12,0x40),2700000); techlines(s,NAVY)
rect(s,0,0,10,0.09,fill=BLUE); rect(s,0,5.535,10,0.09,fill=MAG)
txt(s,0.7,1.0,8.6,0.4,[[('NEXT STEP · 下一步',EN,12,True,CYAN)]])
txt(s,0.7,1.5,8.6,0.9,[[('先来一场 60 分钟的现场样机演示',CN,28,True,WHITE)]])
txt(s,0.7,2.6,8.6,0.7,[[('不必先谈定点。让华翔团队亲眼看一台 L3 全功能可交互样机能做什么',CN,13,False,FOG)],
                       [('—— 这是最轻的一步,也是最有说服力的一步。',CN,13,False,FOG)]],ls=1.3)
cc=rect(s,0.7,3.55,8.6,1.05,fill=RGBColor(0x10,0x22,0x4A),line=RGBColor(0x2B,0x3A,0x66),sh=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08); shadow(cc,alpha=72000); rect(s,0.7,3.63,0.06,0.9,fill=CYAN)
txt(s,0.98,3.68,8.1,0.35,[[('联系我们',CN,11,True,CYAN)]])
txt(s,0.98,4.05,8.1,0.3,[[('严 艇 Tim  17740838008 ｜ tim@jomecdesign.com      李成龙 Jason  15000763365 ｜ jason@jomecdesign.com',CN,9.5,False,FOG)]])
txt(s,0.98,4.35,8.1,0.3,[[('上海卓迈汽车技术有限公司 ｜ www.jomec-auto.com ｜ JOMEC · SHAPING THE FUTURE',CN,9.5,True,WHITE)]])
pg()

out='/home/user/sales-proposal/NICE-智能座舱-售前方案-master-V2.pptx'
prs.save(out)
print('saved',out,'| slides',len(prs.slides._sldIdLst))
